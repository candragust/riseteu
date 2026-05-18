#!/usr/bin/env python3
"""Generate an editable methodology/conceptual-relationship PPTX.

The slide is reconstructed from PowerPoint primitives so the output remains
editable: rounded rectangles, lines, arrows, circles, text boxes, and grouped
vector icons. The source JPEG is only used as a visual reference.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("metodologi_procurement_editable_icons_v2.pptx")

FONT = "Aptos"
TITLE_FONT = "Aptos Display"
ICON_FONT = "Segoe UI Symbol"

NAVY = RGBColor(0x0B, 0x4E, 0x8A)
TITLE_NAVY = RGBColor(0x0D, 0x1B, 0x4C)
TEXT = RGBColor(0x1E, 0x25, 0x3D)
MUTED = RGBColor(0x42, 0x4A, 0x65)
PANEL_FILL = RGBColor(0xFB, 0xFC, 0xFE)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_FILL = RGBColor(0xF5, 0xFB, 0xF0)
PURPLE = RGBColor(0x5E, 0x35, 0xB1)
PURPLE_FILL = RGBColor(0xFA, 0xF7, 0xFF)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
ORANGE_FILL = RGBColor(0xFF, 0xF8, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LINE = RGBColor(0xD8, 0xDE, 0xE8)


def inch(value: float) -> int:
    return Inches(value)


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor | None, width: float = 1.0) -> None:
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_round_rect(slide, x, y, w, h, fill, line, radius_adjust: float = 0.12, width: float = 1.1):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h)
    )
    set_fill(shape, fill)
    set_line(shape, line, width)
    if shape.adjustments:
        shape.adjustments[0] = radius_adjust
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color: RGBColor = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = FONT,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    tf.margin_top = inch(0.01)
    tf.margin_bottom = inch(0.01)
    tf.word_wrap = True
    tf.vertical_anchor = valign

    for idx, line in enumerate(text.split("\n")):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        if line_spacing is not None:
            paragraph.line_spacing = line_spacing
        run = paragraph.runs[0]
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_text_runs(
    slide,
    runs: list[tuple[str, bool]],
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color: RGBColor = TEXT,
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    tf.margin_top = inch(0.01)
    tf.margin_bottom = inch(0.01)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    for value, bold in runs:
        run = paragraph.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_line(slide, x1, y1, x2, y2, color=TITLE_NAVY, width: float = 1.4, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_arrow_shape(slide, x, y, w, h, color=TITLE_NAVY):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inch(x), inch(y), inch(w), inch(h))
    set_fill(arrow, color)
    set_line(arrow, None)
    return arrow


def add_up_arrow_marker(slide, cx: float, y: float, color=TITLE_NAVY):
    marker = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, inch(cx - 0.065), inch(y), inch(0.13), inch(0.105)
    )
    set_fill(marker, color)
    set_line(marker, None)
    return marker


def style_shape_text(
    shape,
    text: str,
    size: float,
    color: RGBColor = WHITE,
    bold: bool = True,
    font: str = ICON_FONT,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = text
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def group_shapes(slide, shapes, name: str):
    group = slide.shapes.add_group_shape(shapes)
    group.name = name
    return group


def add_number_badge(slide, cx: float, cy: float, d: float, color: RGBColor, label: str):
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx - d / 2), inch(cy - d / 2), inch(d), inch(d))
    set_fill(badge, color)
    set_line(badge, None)
    style_shape_text(badge, label, 13.5, WHITE, True, FONT)
    return badge


def draw_building_icon(slide, cx: float, cy: float, d: float, color: RGBColor):
    shapes = []
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx - d / 2), inch(cy - d / 2), inch(d), inch(d))
    set_fill(circle, color)
    set_line(circle, None)
    shapes.append(circle)

    roof = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        inch(cx - d * 0.28),
        inch(cy - d * 0.25),
        inch(d * 0.56),
        inch(d * 0.18),
    )
    set_fill(roof, WHITE)
    set_line(roof, None)
    shapes.append(roof)

    cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(cx - d * 0.24), inch(cy - d * 0.08), inch(d * 0.48), inch(d * 0.06))
    set_fill(cap, WHITE)
    set_line(cap, None)
    shapes.append(cap)

    for offset in (-0.16, 0.0, 0.16):
        col = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            inch(cx + d * offset - d * 0.035),
            inch(cy - d * 0.01),
            inch(d * 0.07),
            inch(d * 0.22),
        )
        set_fill(col, WHITE)
        set_line(col, None)
        shapes.append(col)

    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(cx - d * 0.27), inch(cy + d * 0.23), inch(d * 0.54), inch(d * 0.07))
    set_fill(base, WHITE)
    set_line(base, None)
    shapes.append(base)
    return group_shapes(slide, shapes, "Icon Drivers - Building")


def draw_globe_icon(slide, cx: float, cy: float, d: float, color: RGBColor):
    shapes = []
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx - d / 2), inch(cy - d / 2), inch(d), inch(d))
    set_fill(circle, color)
    set_line(circle, None)
    shapes.append(circle)

    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx - d * 0.31), inch(cy - d * 0.31), inch(d * 0.62), inch(d * 0.62))
    outer.fill.background()
    set_line(outer, WHITE, 2.0)
    shapes.append(outer)

    for width in (0.25,):
        meridian = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            inch(cx - d * width / 2),
            inch(cy - d * 0.31),
            inch(d * width),
            inch(d * 0.62),
        )
        meridian.fill.background()
        set_line(meridian, WHITE, 1.6)
        shapes.append(meridian)

    shapes.append(add_line(slide, cx - d * 0.30, cy, cx + d * 0.30, cy, WHITE, 1.5))
    shapes.append(add_line(slide, cx - d * 0.25, cy - d * 0.14, cx + d * 0.25, cy - d * 0.14, WHITE, 1.2))
    shapes.append(add_line(slide, cx - d * 0.25, cy + d * 0.14, cx + d * 0.25, cy + d * 0.14, WHITE, 1.2))
    return group_shapes(slide, shapes, "Icon Digitalisasi - Globe")


def draw_clock_icon(slide, cx: float, cy: float, d: float, color: RGBColor):
    shapes = []
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx - d / 2), inch(cy - d / 2), inch(d), inch(d))
    set_fill(circle, color)
    set_line(circle, None)
    shapes.append(circle)

    face = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx - d * 0.28), inch(cy - d * 0.28), inch(d * 0.56), inch(d * 0.56))
    face.fill.background()
    set_line(face, WHITE, 2.4)
    shapes.append(face)
    shapes.append(add_line(slide, cx, cy, cx, cy - d * 0.18, WHITE, 2.0))
    shapes.append(add_line(slide, cx, cy, cx + d * 0.16, cy + d * 0.11, WHITE, 2.0))
    return group_shapes(slide, shapes, "Icon Implementasi - Clock")


def draw_chart_icon(slide, cx: float, cy: float, d: float, color: RGBColor):
    shapes = []
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx - d / 2), inch(cy - d / 2), inch(d), inch(d))
    set_fill(circle, color)
    set_line(circle, None)
    shapes.append(circle)

    for xoff, height in [(-0.19, 0.18), (0.0, 0.27), (0.19, 0.39)]:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            inch(cx + d * xoff - d * 0.045),
            inch(cy + d * 0.25 - d * height),
            inch(d * 0.09),
            inch(d * height),
        )
        set_fill(bar, WHITE)
        set_line(bar, None)
        shapes.append(bar)

    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inch(cx - d * 0.19), inch(cy - d * 0.28), inch(d * 0.53), inch(d * 0.13))
    set_fill(arrow, WHITE)
    set_line(arrow, None)
    arrow.rotation = 320
    shapes.append(arrow)
    return group_shapes(slide, shapes, "Icon Kinerja - Chart")


def draw_people_icon(slide, cx: float, cy: float, d: float, color: RGBColor):
    shapes = []
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx - d / 2), inch(cy - d / 2), inch(d), inch(d))
    set_fill(circle, color)
    set_line(circle, None)
    shapes.append(circle)

    for xoff, yoff, radius in [(-0.18, -0.07, 0.11), (0.18, -0.07, 0.11), (0.0, -0.15, 0.13)]:
        head = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            inch(cx + d * xoff - d * radius),
            inch(cy + d * yoff - d * radius),
            inch(d * radius * 2),
            inch(d * radius * 2),
        )
        set_fill(head, WHITE)
        set_line(head, None)
        shapes.append(head)

    for xoff, yoff, width, height in [(-0.18, 0.17, 0.28, 0.22), (0.18, 0.17, 0.28, 0.22), (0.0, 0.16, 0.34, 0.26)]:
        body = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            inch(cx + d * xoff - d * width / 2),
            inch(cy + d * yoff - d * height / 2),
            inch(d * width),
            inch(d * height),
        )
        set_fill(body, WHITE)
        set_line(body, None)
        shapes.append(body)
    return group_shapes(slide, shapes, "Icon Enablers - People")


def draw_target_icon(slide, cx: float, cy: float, d: float, color: RGBColor):
    shapes = []
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx - d / 2), inch(cy - d / 2), inch(d), inch(d))
    set_fill(circle, color)
    set_line(circle, None)
    shapes.append(circle)

    for radius, width in [(0.33, 2.7), (0.22, 2.2), (0.11, 1.8)]:
        ring = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            inch(cx - d * radius),
            inch(cy - d * radius),
            inch(d * radius * 2),
            inch(d * radius * 2),
        )
        ring.fill.background()
        set_line(ring, WHITE, width)
        shapes.append(ring)

    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inch(cx - d * 0.02), inch(cy - d * 0.34), inch(d * 0.46), inch(d * 0.12))
    set_fill(arrow, WHITE)
    set_line(arrow, None)
    arrow.rotation = 320
    shapes.append(arrow)
    return group_shapes(slide, shapes, "Icon Kesimpulan - Target")


def draw_source_icon(slide, x: float, y: float) -> None:
    shapes = []
    icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(0.28), inch(0.28))
    set_fill(icon, NAVY)
    set_line(icon, None)
    if icon.adjustments:
        icon.adjustments[0] = 0.15
    shapes.append(icon)
    for yy, ww in [(0.08, 0.13), (0.14, 0.13), (0.20, 0.09)]:
        shapes.append(add_line(slide, x + 0.08, y + yy, x + 0.08 + ww, y + yy, WHITE, 1.3))
    group_shapes(slide, shapes, "Icon Sumber - List")


def add_main_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    number_title: str,
    title_color: RGBColor,
    fill: RGBColor,
    border: RGBColor,
    icon_drawer,
    body: str,
) -> None:
    add_round_rect(slide, x, y, w, h, fill, border, radius_adjust=0.08, width=1.1)
    cx = x + w / 2
    icon_drawer(slide, cx, y + 0.40, 0.60, title_color)
    add_text(
        slide,
        number_title,
        x + 0.20,
        y + 0.88,
        w - 0.40,
        0.43,
        13.3,
        title_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, x + 0.54, y + 1.34, x + w - 0.54, y + 1.34, title_color, 1.2)
    add_text(
        slide,
        body,
        x + 0.22,
        y + 1.55,
        w - 0.44,
        h - 1.68,
        9.8,
        TEXT,
        align=PP_ALIGN.CENTER,
        line_spacing=0.92,
    )


def add_bottom_card(slide, x, y, w, h, number, color, text):
    add_round_rect(slide, x, y, w, h, WHITE, GRAY_LINE, radius_adjust=0.08, width=1.0)
    add_number_badge(slide, x + 0.30, y + 0.45, 0.38, color, str(number))
    add_text(slide, text, x + 0.63, y + 0.22, w - 0.78, h - 0.25, 10.1, TEXT, line_spacing=0.95)


def build_slide() -> Presentation:
    prs = Presentation()
    prs.slide_width = inch(13.333)
    prs.slide_height = inch(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Soft lower background waves.
    for x, y, w, h, color in [
        (-2.30, 6.05, 5.30, 1.52, RGBColor(0xF0, 0xF5, 0xFB)),
        (-1.95, 6.25, 4.75, 1.18, RGBColor(0xE7, 0xEF, 0xF8)),
        (8.78, 5.72, 5.30, 1.90, RGBColor(0xF0, 0xF5, 0xFB)),
        (9.12, 6.08, 4.95, 1.52, RGBColor(0xE7, 0xEF, 0xF8)),
    ]:
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x), inch(y), inch(w), inch(h))
        set_fill(oval, color)
        set_line(oval, None)

    # Title block.
    add_line(slide, 0.39, 0.25, 0.39, 1.22, NAVY, 4.0)
    add_text(
        slide,
        "Ringkasan Hubungan Konseptual",
        0.73,
        0.31,
        8.9,
        0.46,
        31,
        TITLE_NAVY,
        True,
        font=TITLE_FONT,
    )
    add_text(
        slide,
        "Alur kausal digitalisasi e-procurement terhadap implementasi dan kinerja procurement",
        0.74,
        1.01,
        9.2,
        0.32,
        14.5,
        MUTED,
    )

    # Main process cards.
    card_y, card_w, card_h = 1.52, 2.02, 2.53
    card_xs = [0.48, 3.00, 5.58, 8.16]
    cards = [
        (
            "1. Drivers",
            NAVY,
            PANEL_FILL,
            NAVY,
            draw_building_icon,
            "Faktor pendorong utama\nmendorong organisasi\nmengadopsi e-procurement.",
        ),
        (
            "2. Digitalisasi\ne-Procurement",
            GREEN,
            GREEN_FILL,
            GREEN,
            draw_globe_icon,
            "Sistem dan teknologi digital\nmengubah cara\nprocurement\ndiimplementasikan.",
        ),
        (
            "3. Implementasi\ne-Procurement",
            PURPLE,
            PURPLE_FILL,
            PURPLE,
            draw_clock_icon,
            "Implementasi yang\nlebih baik\nmemperbaiki proses.",
        ),
        (
            "4. Kinerja\nProcurement",
            ORANGE,
            ORANGE_FILL,
            ORANGE,
            draw_chart_icon,
            "Perbaikan proses\nmenghasilkan\npeningkatan kinerja\nprocurement.",
        ),
    ]
    for x, card in zip(card_xs, cards):
        add_main_card(slide, x, card_y, card_w, card_h, *card)

    for x in (2.60, 5.18, 7.75):
        add_arrow_shape(slide, x, 2.53, 0.34, 0.25, NAVY)

    # Enabler bar with dashed causal links.
    enabler_x, enabler_y, enabler_w, enabler_h = 0.48, 4.45, 9.62, 0.90
    add_round_rect(slide, enabler_x, enabler_y, enabler_w, enabler_h, PANEL_FILL, NAVY, radius_adjust=0.07, width=1.0)
    for cx in [x + card_w / 2 for x in card_xs]:
        add_line(slide, cx, enabler_y, cx, card_y + card_h + 0.02, NAVY, 1.1, dashed=True)
        add_up_arrow_marker(slide, cx, card_y + card_h + 0.02, NAVY)
    draw_people_icon(slide, enabler_x + 1.02, enabler_y + 0.45, 0.62, NAVY)
    add_text(slide, "5. Enablers Keberhasilan", enabler_x + 1.50, enabler_y + 0.14, 3.05, 0.22, 12.2, NAVY, True)
    add_text(
        slide,
        "Dukungan manajemen puncak, kompetensi SDM & manajemen perubahan, proses bisnis & tata kelola yang baik,\n"
        "integrasi sistem & infrastruktur TI, serta keamanan, kepercayaan & manajemen risiko.",
        enabler_x + 1.50,
        enabler_y + 0.43,
        enabler_w - 1.72,
        0.39,
        9.7,
        TEXT,
        line_spacing=0.90,
    )

    # Transition chevron and conclusion panel.
    chevron = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, inch(10.28), inch(2.54), inch(0.38), inch(1.76))
    set_fill(chevron, RGBColor(0xE7, 0xF0, 0xFB))
    set_line(chevron, None)

    concl_x, concl_y, concl_w, concl_h = 10.65, 1.98, 2.35, 3.34
    add_round_rect(slide, concl_x, concl_y, concl_w, concl_h, PANEL_FILL, NAVY, radius_adjust=0.08, width=1.1)
    draw_target_icon(slide, concl_x + concl_w / 2, concl_y + 0.58, 0.80, NAVY)
    add_text(slide, "Kesimpulan:", concl_x + 0.22, concl_y + 1.15, concl_w - 0.44, 0.26, 13.6, NAVY, True, PP_ALIGN.CENTER)
    add_line(slide, concl_x + 0.34, concl_y + 1.36, concl_x + concl_w - 0.34, concl_y + 1.36, NAVY, 1.1)
    add_text(
        slide,
        "digitalisasi e-procurement\nberfungsi sebagai\nenabler strategis yang\nmemperkuat implementasi\n"
        "dan meningkatkan\nkinerja procurement\nsecara berkelanjutan.",
        concl_x + 0.23,
        concl_y + 1.60,
        concl_w - 0.46,
        1.42,
        10.1,
        TEXT,
        align=PP_ALIGN.CENTER,
        line_spacing=0.88,
    )

    # Bottom summary cards.
    bottom_y, bottom_h = 5.55, 1.12
    add_bottom_card(slide, 0.47, bottom_y, 2.28, bottom_h, 1, NAVY, "Drivers mendorong\norganisasi mengadopsi\ne-procurement.")
    add_bottom_card(
        slide,
        2.90,
        bottom_y,
        2.38,
        bottom_h,
        2,
        GREEN,
        "Sistem dan teknologi\ndigital mengubah cara\nprocurement\ndiimplementasikan.",
    )
    add_bottom_card(slide, 5.43, bottom_y, 2.30, bottom_h, 3, PURPLE, "Implementasi yang\nlebih baik\nmemperbaiki proses.")
    add_bottom_card(
        slide,
        7.88,
        bottom_y,
        2.22,
        bottom_h,
        4,
        ORANGE,
        "Perbaikan proses\nmenghasilkan\npeningkatan kinerja\nprocurement.",
    )

    # Footer.
    draw_source_icon(slide, 0.56, 6.97)
    add_text_runs(
        slide,
        [
            ("Sumber:  ", True),
            (
                "Quesada et al. (2010); Ramayah et al. (2007); Hsin Chang et al. (2013); "
                "Taherdoost (2023); Jahani et al. (2021); PwC (2024).",
                False,
            ),
        ],
        1.02,
        7.01,
        9.60,
        0.24,
        8.9,
        TEXT,
    )
    return prs


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Generate editable PPTX methodology diagram from PowerPoint shapes."
    )
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=OUT,
        help=f"Output .pptx path, default: {OUT}",
    )
    args = parser.parse_args()
    prs = build_slide()
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(f"Saved: {args.output.resolve()}")


if __name__ == "__main__":
    main()
