#!/usr/bin/env python3
"""Rebuild the image-only procurement deck as editable PowerPoint shapes.

The source deck `digital_procure_diagram_convert.pptx` contains one raster
image per slide. This generator reconstructs the seven slides with native
PowerPoint primitives: text boxes, rounded rectangles, connectors, arrows,
ovals, and simple icon shapes. No raster image is embedded in the output.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
PPTX_OUT = OUT_DIR / "digital_procure_diagram_editable_shapes.pptx"
MANIFEST_OUT = OUT_DIR / "digital_procure_diagram_editable_shapes_manifest.md"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Aptos"
TITLE_FONT = "Aptos Display"

COLORS = {
    "navy": "0B4EA2",
    "title": "06154A",
    "ink": "172033",
    "muted": "3F4E70",
    "line": "D1DAEA",
    "white": "FFFFFF",
    "blue_fill": "F3F8FF",
    "blue_soft": "E8F2FF",
    "green": "2E7D32",
    "green_dark": "0C4A12",
    "green_fill": "F5FBF0",
    "purple": "5E35B1",
    "purple_dark": "160A3B",
    "purple_fill": "F8F4FF",
    "orange": "EF7D00",
    "orange_dark": "C85000",
    "orange_fill": "FFF7ED",
    "gray": "E5EAF3",
}


def inch(value: float) -> int:
    return Inches(value)


def rgb(name: str) -> RGBColor:
    value = COLORS.get(name, name).strip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:], 16))


def set_fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str | None, width: float = 1.0, dashed: bool = False) -> None:
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color: str = "ink",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = FONT,
) -> None:
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    tf.margin_top = inch(0.01)
    tf.margin_bottom = inch(0.01)
    for idx, line in enumerate(text.split("\n")):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        if paragraph.runs:
            run = paragraph.runs[0]
        else:
            run = paragraph.add_run()
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_text_runs(
    slide,
    runs: list[tuple[str, bool]],
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color: str = "ink",
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    tf.margin_top = inch(0.01)
    tf.margin_bottom = inch(0.01)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    for value, bold in runs:
        run = paragraph.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_round_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    line: str | None,
    radius: float = 0.08,
    width: float = 1.0,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h)
    )
    set_fill(shape, fill)
    set_line(shape, line, width)
    if shape.adjustments:
        shape.adjustments[0] = radius
    return shape


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    set_fill(shape, fill)
    set_line(shape, line)
    return shape


def add_oval(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x), inch(y), inch(w), inch(h))
    set_fill(shape, fill)
    set_line(shape, line)
    return shape


def style_shape_text(
    shape,
    text: str,
    size: float,
    color: str = "white",
    bold: bool = True,
    font: str = FONT,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = inch(0.01)
    tf.margin_right = inch(0.01)
    tf.margin_top = inch(0.01)
    tf.margin_bottom = inch(0.01)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = text
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def draw_icon_marks(slide, cx: float, cy: float, d: float, label: str, color: str) -> bool:
    """Draw a compact editable icon inside an existing icon circle."""
    token = label.upper()
    fg = color
    s = d

    def rect(rx, ry, rw, rh, fill_color=fg, line_color=None):
        return add_rect(slide, cx + s * rx, cy + s * ry, s * rw, s * rh, fill_color, line_color)

    def oval(rx, ry, rw, rh, fill_color=fg, line_color=None):
        return add_oval(slide, cx + s * rx, cy + s * ry, s * rw, s * rh, fill_color, line_color)

    def line(x1, y1, x2, y2, width=1.3):
        return add_connector(slide, cx + s * x1, cy + s * y1, cx + s * x2, cy + s * y2, fg, width)

    if token in {"BANK", "GOV", "CORP"}:
        if token == "GOV":
            shield = slide.shapes.add_shape(
                MSO_SHAPE.PENTAGON,
                inch(cx - s * 0.30),
                inch(cy - s * 0.23),
                inch(s * 0.28),
                inch(s * 0.36),
            )
            set_fill(shield, fg)
            set_line(shield, None)
            line(-0.21, -0.02, -0.15, 0.06, 1.4)
            line(-0.15, 0.06, -0.04, -0.12, 1.4)
            clip = rect(0.00, -0.24, 0.26, 0.45, "white", fg)
            clip.fill.background()
            set_line(clip, fg, 1.3)
            rect(0.08, -0.30, 0.10, 0.08)
            line(0.06, -0.10, 0.19, -0.10, 1.0)
            line(0.06, 0.02, 0.19, 0.02, 1.0)
            line(0.06, 0.14, 0.15, 0.14, 1.0)
            return True
        roof = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            inch(cx - s * 0.25),
            inch(cy - s * 0.24),
            inch(s * 0.50),
            inch(s * 0.20),
        )
        set_fill(roof, fg)
        set_line(roof, None)
        rect(-0.26, -0.05, 0.52, 0.06)
        for off in (-0.17, 0.0, 0.17):
            rect(off - 0.035, 0.02, 0.07, 0.24)
        rect(-0.29, 0.28, 0.58, 0.07)
        return True

    if token == "WEB":
        globe = oval(-0.26, -0.26, 0.52, 0.52, "white", fg)
        globe.fill.background()
        set_line(globe, fg, 1.6)
        meridian = oval(-0.11, -0.26, 0.22, 0.52, "white", fg)
        meridian.fill.background()
        set_line(meridian, fg, 1.2)
        line(-0.25, 0.00, 0.25, 0.00, 1.2)
        line(-0.21, -0.13, 0.21, -0.13, 1.0)
        line(-0.21, 0.13, 0.21, 0.13, 1.0)
        return True

    if token == "TIME":
        face = oval(-0.26, -0.26, 0.52, 0.52, "white", fg)
        face.fill.background()
        set_line(face, fg, 1.7)
        line(0, 0, 0, -0.16, 1.5)
        line(0, 0, 0.14, 0.10, 1.5)
        return True

    if token in {"UP", "BAR", "BEN"}:
        for off, height in [(-0.20, 0.20), (-0.03, 0.31), (0.14, 0.43)]:
            rect(off, 0.23 - height, 0.10, height)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            inch(cx - s * 0.13),
            inch(cy - s * 0.26),
            inch(s * 0.42),
            inch(s * 0.12),
        )
        set_fill(arrow, fg)
        set_line(arrow, None)
        arrow.rotation = 320
        if token == "BEN":
            lens = oval(0.05, -0.07, 0.24, 0.24, "white", fg)
            lens.fill.background()
            set_line(lens, fg, 1.4)
            line(0.22, 0.12, 0.34, 0.25, 1.4)
        return True

    if token in {"TEAM", "SDM"}:
        for xoff, yoff, radius in [(-0.17, -0.08, 0.08), (0.17, -0.08, 0.08), (0.0, -0.16, 0.10)]:
            oval(xoff - radius, yoff - radius, radius * 2, radius * 2)
        for xoff, yoff, width, height in [(-0.17, 0.14, 0.23, 0.17), (0.17, 0.14, 0.23, 0.17), (0.0, 0.13, 0.29, 0.21)]:
            oval(xoff - width / 2, yoff - height / 2, width, height)
        if token == "SDM":
            gear = oval(0.18, 0.10, 0.22, 0.22, "white", fg)
            gear.fill.background()
            set_line(gear, fg, 1.3)
            oval(0.245, 0.165, 0.09, 0.09, fg, None)
        return True

    if token == "SPD":
        gauge = slide.shapes.add_shape(
            MSO_SHAPE.ARC,
            inch(cx - s * 0.28),
            inch(cy - s * 0.18),
            inch(s * 0.56),
            inch(s * 0.42),
        )
        set_line(gauge, fg, 1.7)
        line(0, 0.06, 0.20, -0.14, 1.8)
        oval(-0.04, 0.02, 0.08, 0.08, fg, None)
        for xoff, yoff in [(-0.22, 0.04), (-0.12, -0.12), (0.12, -0.12), (0.22, 0.04)]:
            line(xoff, yoff, xoff * 0.86, yoff * 0.86, 1.0)
        return True

    if token == "HND":
        left = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            inch(cx - s * 0.31),
            inch(cy - s * 0.08),
            inch(s * 0.26),
            inch(s * 0.18),
        )
        right = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            inch(cx + s * 0.05),
            inch(cy - s * 0.08),
            inch(s * 0.26),
            inch(s * 0.18),
        )
        for hand in (left, right):
            set_fill(hand, fg)
            set_line(hand, None)
        right.rotation = 180
        line(-0.12, -0.04, 0.12, 0.10, 2.0)
        line(-0.05, 0.05, 0.13, -0.08, 2.0)
        for xoff in (-0.17, -0.10, 0.10, 0.17):
            line(xoff, 0.08, xoff * 0.75, 0.21, 1.1)
        return True

    if token == "!":
        warn = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            inch(cx - s * 0.25),
            inch(cy - s * 0.25),
            inch(s * 0.50),
            inch(s * 0.50),
        )
        set_fill(warn, fg)
        set_line(warn, None)
        add_text(slide, "!", cx - s * 0.06, cy - s * 0.17, s * 0.12, s * 0.26, 10, "white", True, PP_ALIGN.CENTER)
        return True

    if token == "DOC":
        page = rect(-0.18, -0.27, 0.36, 0.54, "white", fg)
        page.fill.background()
        set_line(page, fg, 1.3)
        line(-0.10, -0.06, 0.10, -0.06, 1.1)
        line(-0.10, 0.06, 0.10, 0.06, 1.1)
        line(-0.10, 0.18, 0.05, 0.18, 1.1)
        return True

    if token == "FOL":
        rect(-0.27, -0.10, 0.54, 0.32, fg, None)
        rect(-0.25, -0.20, 0.28, 0.12, fg, None)
        return True

    if token == "MSG":
        box = rect(-0.28, -0.18, 0.56, 0.36, "white", fg)
        box.fill.background()
        set_line(box, fg, 1.3)
        line(-0.27, -0.17, 0.0, 0.04, 1.0)
        line(0.27, -0.17, 0.0, 0.04, 1.0)
        return True

    if token in {"TRK", "FLT"}:
        if token == "FLT":
            funnel = slide.shapes.add_shape(
                MSO_SHAPE.TRAPEZOID,
                inch(cx - s * 0.24),
                inch(cy - s * 0.24),
                inch(s * 0.48),
                inch(s * 0.26),
            )
            set_fill(funnel, fg)
            set_line(funnel, None)
            funnel.rotation = 180
            rect(-0.05, 0.02, 0.10, 0.26)
            oval(-0.03, 0.25, 0.06, 0.06)
            return True
        rect(-0.29, -0.10, 0.34, 0.26)
        rect(0.05, -0.04, 0.20, 0.20)
        oval(-0.22, 0.17, 0.12, 0.12)
        oval(0.11, 0.17, 0.12, 0.12)
        return True

    if token == "PAY":
        card = rect(-0.28, -0.17, 0.56, 0.34, "white", fg)
        card.fill.background()
        set_line(card, fg, 1.3)
        rect(-0.24, -0.10, 0.48, 0.06)
        rect(-0.20, 0.05, 0.13, 0.04)
        rect(0.00, 0.05, 0.16, 0.04)
        return True

    if token in {"BI", "DB"}:
        can = slide.shapes.add_shape(
            MSO_SHAPE.CAN,
            inch(cx - s * 0.24),
            inch(cy - s * 0.26),
            inch(s * 0.48),
            inch(s * 0.52),
        )
        set_fill(can, fg)
        set_line(can, None)
        return True

    if token == "SEC":
        shield = slide.shapes.add_shape(
            MSO_SHAPE.PENTAGON,
            inch(cx - s * 0.23),
            inch(cy - s * 0.25),
            inch(s * 0.46),
            inch(s * 0.48),
        )
        set_fill(shield, fg)
        set_line(shield, None)
        lock = rect(-0.08, -0.02, 0.16, 0.17, "white", None)
        shackle = oval(-0.10, -0.15, 0.20, 0.20, "white", fg)
        shackle.fill.background()
        set_line(shackle, "white", 1.5)
        return True

    if token == "IDEA":
        bulb = oval(-0.15, -0.24, 0.30, 0.32, "white", fg)
        bulb.fill.background()
        set_line(bulb, fg, 1.4)
        rect(-0.10, 0.08, 0.20, 0.08)
        rect(-0.07, 0.18, 0.14, 0.06)
        for angle_x, angle_y in [(-0.28, -0.18), (0.28, -0.18), (-0.32, 0.02), (0.32, 0.02), (0, -0.35)]:
            line(angle_x * 0.72, angle_y * 0.72, angle_x, angle_y, 1.0)
        return True

    if token == "MON":
        monitor = rect(-0.24, -0.18, 0.48, 0.34, "white", fg)
        monitor.fill.background()
        set_line(monitor, fg, 1.4)
        line(-0.08, 0.25, 0.08, 0.25, 1.3)
        line(0, 0.16, 0, 0.25, 1.3)
        line(-0.08, -0.01, -0.01, 0.07, 1.5)
        line(-0.01, 0.07, 0.12, -0.08, 1.5)
        return True

    if token in {"EYE"}:
        eye = oval(-0.28, -0.14, 0.56, 0.28, fg, None)
        oval(-0.10, -0.10, 0.20, 0.20, "white", None)
        return True

    if token == "CHK":
        clip = rect(-0.20, -0.27, 0.40, 0.54, "white", fg)
        clip.fill.background()
        set_line(clip, fg, 1.4)
        rect(-0.10, -0.31, 0.20, 0.10)
        line(-0.10, -0.02, -0.05, 0.04, 1.2)
        line(-0.05, 0.04, 0.08, -0.10, 1.2)
        line(-0.10, 0.15, 0.10, 0.15, 1.0)
        return True

    if token in {"GEAR"}:
        for ang in range(8):
            import math

            a = math.pi * 2 * ang / 8
            rect(math.cos(a) * 0.20 - 0.035, math.sin(a) * 0.20 - 0.035, 0.07, 0.07)
        oval(-0.20, -0.20, 0.40, 0.40)
        center = oval(-0.08, -0.08, 0.16, 0.16, "white", fg)
        center.fill.background()
        set_line(center, fg, 1.0)
        return True

    if token in {"IT", "REP", "DSS"}:
        mon = rect(-0.29, -0.18, 0.38, 0.31, "white", fg)
        mon.fill.background()
        set_line(mon, fg, 1.2)
        if token == "REP":
            rect(-0.22, 0.00, 0.05, 0.08)
            rect(-0.14, -0.07, 0.05, 0.15)
            rect(-0.06, -0.14, 0.05, 0.22)
            lens = oval(0.03, -0.02, 0.22, 0.22, "white", fg)
            lens.fill.background()
            set_line(lens, fg, 1.3)
            line(0.20, 0.15, 0.31, 0.27, 1.3)
        elif token == "DSS":
            oval(0.02, -0.11, 0.24, 0.24, "white", fg)
            rect(0.13, 0.05, 0.12, 0.22)
        else:
            can = slide.shapes.add_shape(
                MSO_SHAPE.CAN,
                inch(cx + s * 0.06),
                inch(cy - s * 0.06),
                inch(s * 0.24),
                inch(s * 0.34),
            )
            set_fill(can, fg)
            set_line(can, None)
        return True

    if token == "AWD":
        oval(-0.17, -0.24, 0.34, 0.34)
        left = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            inch(cx - s * 0.16),
            inch(cy + s * 0.05),
            inch(s * 0.14),
            inch(s * 0.25),
        )
        right = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            inch(cx + s * 0.02),
            inch(cy + s * 0.05),
            inch(s * 0.14),
            inch(s * 0.25),
        )
        for shape in (left, right):
            set_fill(shape, fg)
            set_line(shape, None)
        return True

    if token in {"WIN", "GOAL"}:
        if token == "GOAL":
            for radius in (0.25, 0.16, 0.07):
                ring = oval(-radius, -radius, radius * 2, radius * 2, "white", fg)
                ring.fill.background()
                set_line(ring, fg, 1.4)
            line(0.08, -0.22, 0.27, -0.39, 1.2)
            return True
        cup = rect(-0.15, -0.14, 0.30, 0.25, fg, None)
        bowl = oval(-0.18, -0.27, 0.36, 0.26, fg, None)
        stem = rect(-0.05, 0.10, 0.10, 0.16, fg, None)
        base = rect(-0.18, 0.25, 0.36, 0.06, fg, None)
        for side in (-1, 1):
            handle = oval(side * 0.10 - 0.20 if side < 0 else side * 0.10, -0.17, 0.20, 0.16, "white", fg)
            handle.fill.background()
            set_line(handle, fg, 1.1)
        return bool(cup and bowl and stem and base)

    return False


def add_icon_circle(
    slide,
    cx: float,
    cy: float,
    d: float,
    color: str,
    label: str,
    label_size: float = 15,
    text_color: str = "white",
):
    start_idx = len(slide.shapes)
    circle = add_oval(slide, cx - d / 2, cy - d / 2, d, d, color, None)
    if not draw_icon_marks(slide, cx, cy, d, label, text_color):
        style_shape_text(circle, label, label_size, text_color, True)
        return circle
    icon_shapes = list(slide.shapes)[start_idx:]
    group = slide.shapes.add_group_shape(icon_shapes)
    group.name = f"editable-icon-{label}"
    return group


def add_arrow_head(connector) -> None:
    ln = connector._element.spPr.get_or_add_ln()
    head_end = OxmlElement("a:headEnd")
    head_end.set("type", "triangle")
    head_end.set("w", "med")
    head_end.set("len", "med")
    ln.append(head_end)


def add_connector(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = "navy",
    width: float = 1.4,
    arrow: bool = False,
    dashed: bool = False,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2)
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    if dashed:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        add_arrow_head(connector)
    return connector


def add_right_arrow(slide, x: float, y: float, w: float, h: float, color: str = "navy"):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inch(x), inch(y), inch(w), inch(h))
    set_fill(arrow, color)
    set_line(arrow, None)
    return arrow


def add_title(slide, title: str, subtitle: str, color: str = "title") -> None:
    add_round_rect(slide, 0.43, 0.24, 0.08, 0.94, color, None, radius=0.20)
    add_text(slide, title, 0.78, 0.28, 11.0, 0.80, 31, color, True, font=TITLE_FONT)
    add_text(slide, subtitle, 0.80, 1.10, 10.5, 0.30, 15.5, "muted", True)


def add_dots(slide, color: str, x: float = 12.06, y: float = 0.33) -> None:
    for row in range(4):
        for col in range(7):
            add_oval(slide, x + col * 0.15, y + row * 0.15, 0.025, 0.025, color, None)


def add_waves(slide, color: str = "blue_soft") -> None:
    for x, y, w, h, fill in [
        (-2.4, 6.2, 5.4, 1.5, color),
        (-2.0, 6.4, 5.0, 1.1, "F1F6FD"),
        (10.2, 5.55, 4.8, 2.0, color),
        (9.9, 6.05, 5.2, 1.5, "F1F6FD"),
    ]:
        add_oval(slide, x, y, w, h, fill, None)


def add_footer_source(slide, text: str, color: str = "navy") -> None:
    add_connector(slide, 0.62, 6.86, 10.10, 6.86, "gray", 0.8)
    icon = add_round_rect(slide, 0.75, 6.96, 0.24, 0.24, color, None, radius=0.08)
    style_shape_text(icon, "S", 10, "white", True)
    add_text_runs(slide, [("Sumber:  ", True), (text, False)], 1.17, 6.99, 9.9, 0.25, 8.2)


def add_card_title_text(slide, shape, title: str, body: str, color: str, title_size: float = 14):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = inch(0.12)
    tf.margin_right = inch(0.12)
    tf.margin_top = inch(0.10)
    tf.margin_bottom = inch(0.08)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = title
    p.space_after = Pt(8)
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = rgb(color)

    for line in body.split("\n"):
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.text = line
        p.space_after = Pt(0)
        run = p.runs[0]
        run.font.name = FONT
        run.font.size = Pt(10.2)
        run.font.color.rgb = rgb("ink")


def add_main_flow_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    color: str,
    fill: str,
    icon: str,
) -> None:
    add_round_rect(slide, x, y, w, h, fill, color, radius=0.07, width=1.1)
    add_rect(slide, x + 0.02, y + 0.94, w - 0.04, 0.02, "white", None)
    add_icon_circle(slide, x + w / 2, y + 0.48, 0.86, color, icon, 13)
    add_text(slide, title, x + 0.16, y + 1.02, w - 0.32, 0.78, 16.0, color, True, PP_ALIGN.CENTER)
    add_connector(slide, x + 0.83, y + 1.86, x + w - 0.83, y + 1.86, color, 2.0)
    add_text(slide, body, x + 0.18, y + 2.02, w - 0.36, h - 2.12, 11.4, "ink", True, PP_ALIGN.CENTER)


def slide_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("white")
    add_waves(slide)
    add_dots(slide, "A9C5EA")
    add_title(
        slide,
        "Digitalisasi e-Procurement Mendukung\nImplementasi dan Kinerja Procurement",
        "Transformasi diagram konseptual ke format presentasi",
    )
    xs = [0.76, 3.78, 6.82, 9.91]
    cards = [
        (
            "1. Drivers\nDigitalisasi\ne-Procurement",
            "Faktor pendorong utama\nyang mendorong adopsi\ndan transformasi\ne-procurement.",
            "navy",
            "blue_fill",
            "BANK",
        ),
        (
            "2. Digitalisasi\ne-Procurement\n(Sistem & Teknologi)",
            "Pemanfaatan sistem dan\nteknologi digital untuk\nmendukung proses\ne-procurement.",
            "green",
            "green_fill",
            "WEB",
        ),
        (
            "3. Implementasi\ne-Procurement\n(Dampak pada Proses)",
            "Dampak penerapan\ne-procurement terhadap\nproses: lebih cepat,\ntransparan, akuntabel,\ndan berkualitas.",
            "purple",
            "purple_fill",
            "TIME",
        ),
        (
            "4. Kinerja\nProcurement\n(Hasil Kinerja)",
            "Peningkatan hasil kinerja\nprocurement: efisiensi\nbiaya, kualitas, ketepatan\nwaktu, dan nilai organisasi.",
            "orange",
            "orange_fill",
            "UP",
        ),
    ]
    for x, card in zip(xs, cards):
        add_main_flow_card(slide, x, 2.08, 2.56, 3.35, *card)
    for x in (3.40, 6.45, 9.50):
        add_right_arrow(slide, x, 3.42, 0.36, 0.24, "navy")

    bar = add_round_rect(slide, 0.76, 5.63, 11.72, 1.05, "blue_fill", "navy", radius=0.05, width=1.1)
    add_icon_circle(slide, 1.72, 6.15, 0.78, "navy", "TEAM", 10)
    add_text(slide, "5. Enablers Keberhasilan", 2.62, 5.80, 3.70, 0.30, 16.2, "navy", True)
    add_text(
        slide,
        "Dukungan manajemen puncak, kompetensi SDM & manajemen perubahan, proses bisnis & tata kelola yang baik,\n"
        "integrasi sistem & infrastruktur TI, serta keamanan, kepercayaan & manajemen risiko.",
        2.62,
        6.14,
        8.92,
        0.36,
        10.9,
        "ink",
        True,
    )
    bar.name = "editable-enablers-bar"
    add_footer_source(
        slide,
        "Quesada et al. (2010); Ramayah et al. (2007); Hsin Chang et al. (2013); Taherdoost (2023); Jahani et al. (2021); PwC (2024).",
    )


def slide_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("white")
    add_waves(slide)
    add_dots(slide, "A9C5EA")
    add_title(
        slide,
        "1. Drivers Digitalisasi e-Procurement",
        "Faktor pendorong utama adopsi dan transformasi digital procurement",
    )
    add_round_rect(slide, 0.78, 1.58, 3.55, 4.70, "blue_fill", "navy", radius=0.06, width=1.2)
    add_icon_circle(slide, 2.56, 2.36, 1.28, "navy", "BANK", 12)
    add_text(slide, "1. Drivers\nDigitalisasi\ne-Procurement", 1.10, 3.15, 2.90, 1.10, 18.2, "navy", True, PP_ALIGN.CENTER)
    add_connector(slide, 1.78, 4.45, 3.34, 4.45, "navy", 2.0)
    add_text(
        slide,
        "Faktor pendorong utama\nyang mendorong adopsi\ndan transformasi\ne-procurement.",
        1.20,
        4.75,
        2.70,
        1.05,
        14.2,
        "ink",
        True,
        PP_ALIGN.CENTER,
    )
    add_right_arrow(slide, 4.46, 3.65, 0.30, 0.35, "navy")

    rows = [
        ("BANK", "navy", "1. Tuntutan transparansi, akuntabilitas & kepatuhan", "(Croom & Brandon-Jones, 2007; OECD, 2016)"),
        ("$", "green", "2. Efisiensi biaya & nilai terbaik (value for money)", "(Thai, 2001; Mabert et al., 2003)"),
        ("UP", "orange", "3. Peningkatan kinerja & daya saing organisasi", "(Gupta & Narain, 2015; Alawadhi & Alshurideh, 2024)"),
        ("TEAM", "purple", "4. Kompleksitas rantai pasok & kolaborasi pemasok", "(Monczka et al., 2008; Handfield et al., 2015)"),
    ]
    for idx, (icon, color, title, source) in enumerate(rows):
        y = 1.72 + idx * 1.20
        add_round_rect(slide, 4.95, y, 7.72, 0.95, "white", "line", radius=0.05, width=1.1)
        add_icon_circle(slide, 5.60, y + 0.47, 0.82, color, icon, 12)
        add_text(slide, title, 6.32, y + 0.18, 6.0, 0.30, 15.6, "title", True)
        add_text(slide, source, 6.32, y + 0.57, 5.8, 0.26, 11.6, "ink", True)

    add_round_rect(slide, 0.78, 6.48, 11.88, 0.80, "blue_fill", "A9C5EA", radius=0.06, width=1.0)
    add_icon_circle(slide, 1.36, 6.88, 0.62, "navy", "IDEA", 8)
    add_text_runs(
        slide,
        [
            ("Implikasi:  ", True),
            ("faktor-faktor ini mendorong organisasi untuk mengadopsi e-procurement\nsebagai fondasi peningkatan implementasi procurement.", True),
        ],
        2.00,
        6.64,
        9.9,
        0.50,
        12.0,
        "ink",
    )


def slide_3(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("white")
    add_waves(slide, "E2F1DF")
    add_dots(slide, "89B884")
    add_title(
        slide,
        "2. Digitalisasi e-Procurement (Sistem & Teknologi)",
        "Komponen digital yang mendukung implementasi procurement",
        "green_dark",
    )
    add_round_rect(slide, 0.72, 1.54, 8.30, 5.28, "green_fill", "green", radius=0.05, width=1.1)
    items = [
        ("WEB", "e-Tendering / e-Sourcing", "Tender elektronik, e-auction"),
        ("DOC", "e-Ordering & e-Catalog", "Pemesanan & katalog elektronik"),
        ("TRK", "e-Contract Management", "Manajemen kontrak elektronik"),
        ("PAY", "e-Payment", "Pembayaran elektronik"),
        ("BI", "Data & Analytics", "Event log, dashboard, BI"),
        ("SEC", "Keamanan informasi & integrasi", "Keamanan data, integrasi sistem"),
    ]
    for idx, (icon, title, body) in enumerate(items):
        y = 1.70 + idx * 0.84
        add_round_rect(slide, 1.10, y, 7.60, 0.73, "white", "9DCB91", radius=0.05, width=0.8)
        add_icon_circle(slide, 1.76, y + 0.36, 0.56, "green", icon, 8)
        add_text(slide, title, 2.55, y + 0.12, 5.2, 0.26, 15.3, "green_dark", True)
        add_text(slide, body, 2.55, y + 0.44, 5.8, 0.22, 11.6, "ink", True)

    add_round_rect(slide, 9.55, 2.76, 2.98, 2.40, "green_fill", "green", radius=0.05, width=1.0)
    add_rect(slide, 9.56, 2.76, 2.96, 0.78, "EAF6E5", None)
    add_icon_circle(slide, 10.03, 3.15, 0.64, "green", "IDEA", 8)
    add_text(slide, "Makna", 10.52, 3.00, 1.4, 0.30, 16.0, "green_dark", True)
    add_text(
        slide,
        "Digitalisasi tidak hanya berarti\notomatisasi transaksi, tetapi juga\nintegrasi data, visibilitas proses,\ndan kemampuan analitik.",
        9.80,
        3.75,
        2.30,
        1.05,
        11.5,
        "ink",
        True,
    )
    add_footer_source(
        slide,
        "Diadaptasi dari Quesada et al. (2010); Hsin Chang et al. (2013); Taherdoost (2023); Jahani et al. (2021).",
        "green",
    )


def list_panel(slide, x: float, y: float, w: float, h: float, title: str, icon: str, color: str, items: list[tuple[str, str]]):
    add_round_rect(slide, x, y, w, h, "purple_fill", color, radius=0.06, width=1.0)
    add_icon_circle(slide, x + w / 2, y + 0.48, 0.82, color, icon, 10)
    add_text(slide, title, x + 0.16, y + 1.02, w - 0.32, 0.32, 15.4, color, True, PP_ALIGN.CENTER)
    add_connector(slide, x + 0.26, y + 1.42, x + w - 0.26, y + 1.42, color, 1.2)
    for idx, (small_icon, text) in enumerate(items):
        row_y = y + 1.66 + idx * 0.54
        add_icon_circle(slide, x + 0.40, row_y + 0.16, 0.30, color, small_icon, 5)
        add_text(slide, text, x + 0.72, row_y - 0.01, w - 0.86, 0.38, 10.2, "ink", True)
        if idx < len(items) - 1:
            add_connector(slide, x + 0.22, row_y + 0.43, x + w - 0.22, row_y + 0.43, "gray", 0.5)


def slide_4(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("white")
    add_waves(slide, "EEE4FA")
    add_dots(slide, "9A74D0")
    add_title(
        slide,
        "3. Implementasi e-Procurement\n(Dampak pada Proses)",
        "Perubahan pada pelaksanaan proses procurement setelah digitalisasi",
        "purple_dark",
    )
    before = [
        ("DOC", "Proses manual"),
        ("FOL", "Informasi tersebar\ndi banyak tempat"),
        ("MSG", "Dokumen & komunikasi\ntidak terintegrasi"),
        ("!", "Sulit dilacak &\nrawan kesalahan"),
        ("BAR", "Keputusan berbasis data\nterbatas"),
    ]
    after = [
        ("SPD", "Proses lebih cepat\ndan efisien"),
        ("EYE", "Transparansi & visibilitas\nmeningkat"),
        ("CHK", "Standardisasi & kepatuhan\nmeningkat"),
        ("HND", "Kolaborasi dengan pemasok\nlebih efektif"),
        ("DB", "Kualitas data & informasi\nmeningkat"),
    ]
    list_panel(slide, 0.62, 1.86, 2.52, 4.28, "Sebelum", "BANK", "purple", before)
    add_right_arrow(slide, 3.24, 3.72, 0.70, 0.44, "purple")
    list_panel(slide, 4.02, 1.86, 2.90, 4.28, "Sesudah digitalisasi", "MON", "purple", after)

    outcomes = [
        ("TIME", "1", "Proses lebih cepat\n(cycle time berkurang)", "(Shen et al., 2014)"),
        ("EYE", "2", "Transparansi & visibilitas\nproses meningkat", "(Croom & Brandon-Jones, 2007)"),
        ("CHK", "3", "Standardisasi & kepatuhan\nproses meningkat", "(OECD, 2016)"),
        ("TEAM", "4", "Kolaborasi dengan pemasok\nlebih efektif", "(Handfield et al., 2015)"),
        ("DB", "5", "Kualitas data & informasi\nmeningkat", "(Schonberger, 2016)"),
    ]
    for idx, (icon, num, title, source) in enumerate(outcomes):
        y = 1.76 + idx * 0.91
        add_round_rect(slide, 7.55, y, 5.08, 0.80, "white", "purple", radius=0.04, width=0.9)
        add_icon_circle(slide, 8.15, y + 0.40, 0.66, "purple", icon, 8)
        add_connector(slide, 8.78, y + 0.14, 8.78, y + 0.66, "line", 0.9)
        add_icon_circle(slide, 9.16, y + 0.40, 0.35, "purple", num, 11)
        add_text(slide, title, 9.55, y + 0.10, 2.75, 0.36, 12.8, "ink", True)
        add_text(slide, source, 9.55, y + 0.54, 2.75, 0.22, 9.4, "ink", True)

    add_round_rect(slide, 0.62, 6.43, 12.01, 0.90, "purple", "purple_dark", radius=0.05, width=1.0)
    add_icon_circle(slide, 1.70, 6.88, 0.76, "white", "IDEA", 8, "purple")
    add_text(
        slide,
        "Intinya, implementasi e-procurement memperbaiki cara proses procurement dijalankan\n"
        "dari sisi kecepatan, keterlacakan, konsistensi, dan kualitas informasi.",
        2.35,
        6.60,
        9.75,
        0.52,
        13.8,
        "white",
        True,
    )


def outcome_card(slide, x, y, w, h, icon, title, source, color="orange"):
    add_round_rect(slide, x, y, w, h, "orange_fill", color, radius=0.05, width=1.0)
    add_oval(slide, x + 0.21, y + 0.23, 1.04, 1.04, "FFF3DE", "F7DBB3")
    add_icon_circle(slide, x + 0.73, y + 0.75, 0.76, color, icon, 12)
    add_text(slide, title, x + 1.52, y + 0.28, w - 1.68, 0.68, 14.2, "ink", True)
    add_text(slide, source, x + 1.52, y + 1.00, w - 1.68, 0.46, 11.0, "muted", True)


def slide_5(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("white")
    add_waves(slide, "FDE7C8")
    add_dots(slide, "F59E42")
    add_title(
        slide,
        "4. Kinerja Procurement (Hasil Kinerja)",
        "Outcome utama dari implementasi e-procurement yang efektif",
        "orange_dark",
    )
    outcome_card(slide, 0.53, 1.70, 3.78, 1.88, "$", "1. Penghematan biaya\n(cost savings)", "(Thai, 2001; Gupta &\nNarain, 2015)")
    outcome_card(slide, 4.62, 1.70, 3.82, 1.88, "AWD", "2. Peningkatan kualitas\nbarang/jasa", "(Jama & Mohamud, 2023)")
    outcome_card(slide, 8.76, 1.70, 3.95, 1.88, "TRK", "3. Ketepatan waktu\npengiriman\n(on-time delivery)", "(Gupta & Narain, 2015)")
    outcome_card(slide, 0.55, 3.86, 5.90, 1.76, "HND", "4. Hubungan pemasok\nlebih kuat & berkelanjutan", "(Revilla & Knoppen, 2015;\nSemuel et al., 2018)")
    outcome_card(slide, 6.78, 3.86, 5.91, 1.76, "UP", "5. Keunggulan bersaing &\nkeberlanjutan kinerja", "(Alawadhi & Alshurideh, 2024;\nSalimian et al., 2021)")

    add_round_rect(slide, 1.28, 6.01, 10.70, 0.96, "orange_fill", "orange", radius=0.05, width=1.0)
    add_round_rect(slide, 1.40, 6.09, 1.18, 0.78, "orange", None, radius=0.07)
    add_icon_circle(slide, 1.99, 6.48, 0.62, "orange", "WIN", 9)
    add_text(
        slide,
        "Kinerja procurement meningkat ketika digitalisasi mendukung\nimplementasi proses secara konsisten dan terukur.",
        2.92,
        6.20,
        8.3,
        0.56,
        15.0,
        "orange_dark",
        True,
    )


def enabler_card(slide, x, y, w, h, icon, number, title, source):
    add_round_rect(slide, x, y, w, h, "blue_fill", "navy", radius=0.05, width=1.0)
    add_oval(slide, x + 0.02, y + 0.94, w - 0.04, 0.90, "white", None)
    add_icon_circle(slide, x + w / 2, y + 0.60, 1.08, "navy", icon, 10)
    add_icon_circle(slide, x + w / 2, y + 1.26, 0.30, "D5E9FF", number, 9, "navy")
    add_text(slide, title, x + 0.22, y + 1.56, w - 0.44, 0.66, 14.0, "navy", True, PP_ALIGN.CENTER)
    add_text(slide, source, x + 0.22, y + 2.22, w - 0.44, 0.28, 10.8, "ink", True, PP_ALIGN.CENTER)


def slide_6(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("white")
    add_waves(slide)
    add_dots(slide, "A9C5EA")
    add_title(
        slide,
        "5. Enablers Keberhasilan",
        "Faktor pendukung yang memfasilitasi adopsi, implementasi, dan hasil kinerja e-procurement",
    )
    enabler_card(slide, 0.98, 1.45, 3.62, 2.50, "TEAM", "1", "Dukungan manajemen puncak", "(Thong, 1999)")
    enabler_card(slide, 4.80, 1.45, 3.64, 2.50, "SDM", "2", "Kompetensi SDM &\nmanajemen perubahan", "(Lacity & Hirschheim, 1993)")
    enabler_card(slide, 8.66, 1.45, 3.60, 2.50, "GEAR", "3", "Proses bisnis &\ntata kelola yang baik", "(Nielsen, 2006; OECD, 2016)")
    enabler_card(slide, 2.74, 4.12, 3.82, 2.18, "IT", "4", "Integrasi sistem &\ninfrastruktur TI", "(Monczka et al., 2008)")
    enabler_card(slide, 6.74, 4.12, 3.82, 2.18, "SEC", "5", "Keamanan, kepercayaan &\nmanajemen risiko", "(Cavusoglu et al., 2004)")

    add_round_rect(slide, 1.52, 6.48, 10.24, 0.75, "navy", "navy", radius=0.16, width=1.0)
    add_icon_circle(slide, 2.20, 6.86, 0.62, "white", "SEC", 8, "navy")
    add_text(
        slide,
        "Tanpa enablers yang memadai, manfaat digitalisasi e-procurement sulit diwujudkan secara optimal.",
        2.68,
        6.72,
        8.55,
        0.30,
        13.2,
        "white",
        True,
    )


def contribution_card(slide, x, y, w, h, number, color, icon, title, body):
    add_round_rect(slide, x, y, w, h, "white", color, radius=0.05, width=1.0)
    add_icon_circle(slide, x + 0.27, y + 0.25, 0.32, color, str(number), 10)
    add_icon_circle(slide, x + 1.00, y + h / 2 + 0.05, 0.82, color, icon, 9)
    add_text(slide, title, x + 1.62, y + 0.22, w - 1.82, 0.50, 13.2, "navy", True)
    add_text(slide, body, x + 1.62, y + 0.86, w - 1.82, 0.58, 10.6, "ink", True)


def slide_7(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("white")
    add_waves(slide)
    add_dots(slide, "A9C5EA", 0.17, 0.24)
    add_dots(slide, "A9C5EA", 11.58, 6.42)
    add_text(
        slide,
        "Diagram Konseptual: Kontribusi Praktis bagi Korporasi",
        1.26,
        0.22,
        10.9,
        0.48,
        28.0,
        "title",
        True,
        PP_ALIGN.CENTER,
        font=TITLE_FONT,
    )
    add_text(
        slide,
        "Manfaat praktis penelitian untuk unit procurement dan manajemen korporasi",
        2.42,
        0.82,
        8.6,
        0.28,
        14.5,
        "muted",
        True,
        PP_ALIGN.CENTER,
    )
    contribution_card(
        slide,
        0.45,
        1.20,
        3.75,
        1.62,
        1,
        "navy",
        "REP",
        "Pemahaman Kinerja\nTender Secara Rinci",
        "Memahami kinerja proses dari\nlevel kasus hingga level aktivitas.",
    )
    contribution_card(
        slide,
        4.52,
        1.20,
        3.55,
        1.62,
        2,
        "orange",
        "FLT",
        "Identifikasi Bottleneck\ndan Variasi Proses",
        "Menemukan aktivitas kritis,\nbottleneck, variasi proses, dan\nsumber keterlambatan.",
    )
    contribution_card(
        slide,
        8.35,
        1.20,
        3.86,
        1.62,
        3,
        "green",
        "BEN",
        "Benchmarking\nAntarentitas yang Objektif",
        "Membandingkan kinerja proses\ntender antarunit berdasarkan\nproses aktual.",
    )
    contribution_card(
        slide,
        0.45,
        4.48,
        3.75,
        1.52,
        4,
        "navy",
        "DB",
        "Standardisasi Event\nLog dan Metrik",
        "Mendorong penyelarasan\ndefinisi aktivitas, event log,\ndan metrik kinerja.",
    )
    contribution_card(
        slide,
        4.52,
        4.48,
        3.55,
        1.52,
        5,
        "orange",
        "DSS",
        "Dasar Pengambilan\nKeputusan Berbasis Data",
        "Mendukung keputusan\nmanajerial yang lebih tepat\nberdasarkan bukti proses.",
    )
    contribution_card(
        slide,
        8.35,
        4.48,
        3.86,
        1.52,
        6,
        "green",
        "GOV",
        "Penguatan Tata Kelola\nProcurement",
        "Meningkatkan akuntabilitas,\nauditabilitas, dan pengendalian\nproses tender.",
    )

    add_round_rect(slide, 4.28, 3.04, 3.78, 1.32, "white", "navy", radius=0.08, width=3.0)
    add_icon_circle(slide, 5.02, 3.70, 0.90, "navy", "CORP", 8)
    add_connector(slide, 5.72, 3.18, 5.72, 4.22, "line", 1.0)
    add_text(slide, "Kontribusi\nPraktis bagi\nKorporasi", 6.08, 3.24, 1.65, 0.82, 18.0, "navy", True, PP_ALIGN.CENTER)

    add_connector(slide, 2.30, 2.82, 4.28, 3.70, "navy", 1.2, True)
    add_connector(slide, 6.30, 2.82, 6.30, 3.04, "navy", 1.2, True)
    add_connector(slide, 10.10, 2.82, 8.06, 3.70, "navy", 1.2, True)
    add_connector(slide, 2.30, 4.48, 4.28, 3.70, "navy", 1.2, True)
    add_connector(slide, 6.30, 4.36, 6.30, 4.48, "navy", 1.2, True)
    add_connector(slide, 10.10, 4.48, 8.06, 3.70, "navy", 1.2, True)

    add_round_rect(slide, 1.25, 6.18, 9.70, 0.78, "navy", "navy", radius=0.10)
    add_icon_circle(slide, 1.84, 6.57, 0.64, "white", "GOAL", 7, "navy")
    add_text(
        slide,
        "Arah manfaat: analisis proses -> benchmarking -> keputusan berbasis data ->\n"
        "tata kelola procurement yang lebih akuntabel",
        2.78,
        6.33,
        7.48,
        0.36,
        13.8,
        "white",
        True,
        PP_ALIGN.CENTER,
    )


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    return prs


def count_shapes(prs: Presentation) -> tuple[int, int, int, int]:
    total = pictures = connectors = groups = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            total += 1
            if shape.shape_type == 13:
                pictures += 1
            if shape.shape_type == 9:
                connectors += 1
            if shape.shape_type == 6:
                groups += 1
    return total, pictures, connectors, groups


def write_manifest(prs: Presentation) -> None:
    total, pictures, connectors, groups = count_shapes(prs)
    MANIFEST_OUT.write_text(
        "\n".join(
            [
                "# Manifest Digital Procurement Editable Shapes",
                "",
                f"Dibuat: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
                "",
                "Output:",
                "",
                f"- `{PPTX_OUT.name}`",
                "",
                "Karakter:",
                "",
                "- 7 slide widescreen.",
                "- Direkonstruksi dari deck raster `digital_procure_diagram_convert.pptx`.",
                "- Semua komponen utama berupa shape/line/text PowerPoint editable.",
                "- Tidak memakai gambar raster.",
                f"- Total top-level shapes: {total}.",
                f"- Icon/group objects: {groups}.",
                f"- Pictures: {pictures}.",
                f"- Connectors/lines: {connectors}.",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


def main() -> None:
    prs = build_deck()
    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_OUT)
    write_manifest(prs)
    total, pictures, connectors, groups = count_shapes(prs)
    print(
        f"Saved {PPTX_OUT} slides={len(prs.slides)} "
        f"shapes={total} groups={groups} pictures={pictures} connectors={connectors}"
    )


if __name__ == "__main__":
    main()
